#!/usr/bin/env python3
"""Add Agentic AI slides to the Shehn.AI presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
ROSE_500   = RGBColor(0xD4, 0x73, 0x6E)
ROSE_100   = RGBColor(0xFF, 0xE8, 0xE4)
ROSE_50    = RGBColor(0xFF, 0xF5, 0xF3)
SAGE_600   = RGBColor(0x4A, 0x6B, 0x4A)
GRAY_900   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_800   = RGBColor(0x2D, 0x2D, 0x3F)
GRAY_700   = RGBColor(0x37, 0x41, 0x51)
GRAY_600   = RGBColor(0x4B, 0x55, 0x63)
GRAY_500   = RGBColor(0x6B, 0x72, 0x80)
GRAY_400   = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_200   = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100   = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_50    = RGBColor(0xF9, 0xFA, 0xFB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER_500  = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
BLUE_500   = RGBColor(0x3B, 0x82, 0xF6)
PURPLE_500 = RGBColor(0x8B, 0x5C, 0xF6)
INDIGO     = RGBColor(0x63, 0x66, 0xF1)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
TEAL       = RGBColor(0x14, 0xB8, 0xA6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.08
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox

def add_slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3), f"{num}/{total}", font_size=9, color=GRAY_400, align=PP_ALIGN.RIGHT)

# Load existing presentation
prs = Presentation("/Users/snehaaiyer/Downloads/shehnai-local/Shehnai_Product_Walkthrough.pptx")
TOTAL = 21  # 18 existing + 3 new

# We'll add slides at the end and they'll be the new slides 16, 17, 18 (before closing)
# But python-pptx doesn't support insert, so we add at end

# ══════════════════════════════════════════════════
# SLIDE 19 — SECTION: AGENTIC AI
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRAY_900)
# Decorative
add_rect(slide, Inches(0), Inches(3.2), Inches(1.5), Inches(0.06), ROSE_500)
add_text(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.8), "Agentic AI Architecture", font_size=40, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(0.5), "Multi-agent orchestration with CrewAI — 15 specialized agents across 3 crews", font_size=18, color=GRAY_400)
add_slide_number(slide, 19, TOTAL)

# ══════════════════════════════════════════════════
# SLIDE 20 — CREWAI AGENT SYSTEM
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRAY_900)
add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5), "CrewAI Multi-Agent System", font_size=28, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(0.9), Inches(1.2), Inches(0.04), ROSE_500)
add_text(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.4), "15 specialized agents organized in 3 crews, each with domain expertise and custom tools", font_size=13, color=GRAY_400)

# ── CREW 1: Production Wedding Agents ──
crew1_x = Inches(0.4)
crew1_w = Inches(4.0)
add_shape(slide, crew1_x, Inches(1.6), crew1_w, Inches(5.3), GRAY_800)
add_rect(slide, crew1_x, Inches(1.6), crew1_w, Inches(0.06), ROSE_500)
add_text(slide, crew1_x + Inches(0.2), Inches(1.75), Inches(3.6), Inches(0.3), "Crew 1: Core Planning Engine", font_size=13, color=ROSE_500, bold=True)
add_text(slide, crew1_x + Inches(0.2), Inches(2.05), Inches(3.6), Inches(0.25), "5 agents  ·  Parallel + Sequential  ·  Memory-injected", font_size=9, color=GRAY_500)

agents1 = [
    ("🧮", "Budget Agent", "15+ yrs Indian wedding finance\nAllocates across 6 categories\nCity-tier aware pricing", AMBER_500),
    ("🔍", "Vendor Agent", "Deep vendor industry knowledge\nMatch scoring (0-100)\nTools: query_vendors, get_rates", BLUE_500),
    ("🎨", "Style Agent", "Color theory + cultural expertise\nTheme & decor recommendations\nRegional style knowledge", PURPLE_500),
    ("📅", "Timeline Agent", "Multi-day ceremony sequencing\nAuspicious timing (muhurat)\nVendor scheduling", EMERALD),
    ("💬", "Comms Agent", "Vendor relationship management\nDraft inquiries & negotiations\nBooking confirmations", CYAN),
]
for i, (icon, name, desc, color) in enumerate(agents1):
    y = Inches(2.4 + i * 0.93)
    agent_box = add_shape(slide, crew1_x + Inches(0.15), y, Inches(3.7), Inches(0.82), RGBColor(0x23, 0x23, 0x35))
    add_rect(slide, crew1_x + Inches(0.15), y, Inches(0.06), Inches(0.82), color)
    add_text(slide, crew1_x + Inches(0.35), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, font_size=14, color=WHITE)
    add_text(slide, crew1_x + Inches(0.75), y + Inches(0.05), Inches(2.0), Inches(0.22), name, font_size=10, color=WHITE, bold=True)
    add_text(slide, crew1_x + Inches(0.75), y + Inches(0.28), Inches(2.9), Inches(0.5), desc, font_size=8, color=GRAY_400)

# ── CREW 2: Intelligent Vendor Selection ──
crew2_x = Inches(4.6)
crew2_w = Inches(4.0)
add_shape(slide, crew2_x, Inches(1.6), crew2_w, Inches(5.3), GRAY_800)
add_rect(slide, crew2_x, Inches(1.6), crew2_w, Inches(0.06), BLUE_500)
add_text(slide, crew2_x + Inches(0.2), Inches(1.75), Inches(3.6), Inches(0.3), "Crew 2: Vendor Selection Pipeline", font_size=13, color=BLUE_500, bold=True)
add_text(slide, crew2_x + Inches(0.2), Inches(2.05), Inches(3.6), Inches(0.25), "5 agents  ·  Sequential pipeline  ·  Weighted scoring", font_size=9, color=GRAY_500)

agents2 = [
    ("🏗️", "Filtering Specialist", "Budget/capacity/location filtering\nQualification criteria\nSerperDev web search tool", ORANGE),
    ("📊", "Scoring Specialist", "Multi-factor match algorithm\n7 weighted dimensions (100pts)\nBudget 25, Style 20, Location 15", BLUE_500),
    ("💰", "Budget Validator", "Category budget alignment\nCost-impact analysis\nOverspend/underspend detection", AMBER_500),
    ("🎭", "Style Matcher", "Portfolio ↔ theme alignment\nCultural appropriateness\nAesthetic compatibility", PURPLE_500),
    ("🔬", "Research Specialist", "Credential verification\nMarket positioning analysis\nRisk assessment", TEAL),
]
for i, (icon, name, desc, color) in enumerate(agents2):
    y = Inches(2.4 + i * 0.93)
    agent_box = add_shape(slide, crew2_x + Inches(0.15), y, Inches(3.7), Inches(0.82), RGBColor(0x23, 0x23, 0x35))
    add_rect(slide, crew2_x + Inches(0.15), y, Inches(0.06), Inches(0.82), color)
    add_text(slide, crew2_x + Inches(0.35), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, font_size=14, color=WHITE)
    add_text(slide, crew2_x + Inches(0.75), y + Inches(0.05), Inches(2.0), Inches(0.22), name, font_size=10, color=WHITE, bold=True)
    add_text(slide, crew2_x + Inches(0.75), y + Inches(0.28), Inches(2.9), Inches(0.5), desc, font_size=8, color=GRAY_400)

# ── CREW 3: Indian Wedding Specialized ──
crew3_x = Inches(8.8)
crew3_w = Inches(4.0)
add_shape(slide, crew3_x, Inches(1.6), crew3_w, Inches(5.3), GRAY_800)
add_rect(slide, crew3_x, Inches(1.6), crew3_w, Inches(0.06), EMERALD)
add_text(slide, crew3_x + Inches(0.2), Inches(1.75), Inches(3.6), Inches(0.3), "Crew 3: Indian Wedding Experts", font_size=13, color=EMERALD, bold=True)
add_text(slide, crew3_x + Inches(0.2), Inches(2.05), Inches(3.6), Inches(0.25), "5 agents  ·  Domain-specialized  ·  Cultural context", font_size=9, color=GRAY_500)

agents3 = [
    ("🏛️", "Venue Specialist", "Mandap setup expertise\nGuest capacity planning\nBaraat procession logistics", BLUE_500),
    ("💵", "Budget Expert", "Multi-ceremony cost modeling\nRegional price variations\nLarge guest list optimization", AMBER_500),
    ("🌺", "Decoration Expert", "Regional styles (Rajasthani, South)\nTraditional color symbolism\nMandap designs + floral", ROSE_500),
    ("🍽️", "Catering Specialist", "Regional cuisines (30+ types)\nJain/veg/dietary requirements\nLive counter planning", EMERALD),
    ("⏰", "Timeline Coordinator", "Muhurat-aware scheduling\nMulti-day ceremony flow\nVendor arrival coordination", INDIGO),
]
for i, (icon, name, desc, color) in enumerate(agents3):
    y = Inches(2.4 + i * 0.93)
    agent_box = add_shape(slide, crew3_x + Inches(0.15), y, Inches(3.7), Inches(0.82), RGBColor(0x23, 0x23, 0x35))
    add_rect(slide, crew3_x + Inches(0.15), y, Inches(0.06), Inches(0.82), color)
    add_text(slide, crew3_x + Inches(0.35), y + Inches(0.05), Inches(0.4), Inches(0.3), icon, font_size=14, color=WHITE)
    add_text(slide, crew3_x + Inches(0.75), y + Inches(0.05), Inches(2.0), Inches(0.22), name, font_size=10, color=WHITE, bold=True)
    add_text(slide, crew3_x + Inches(0.75), y + Inches(0.28), Inches(2.9), Inches(0.5), desc, font_size=8, color=GRAY_400)

add_slide_number(slide, 20, TOTAL)

# ══════════════════════════════════════════════════
# SLIDE 21 — ORCHESTRATION PATTERNS + FALLBACK
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ROSE_500)
add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5), "Agent Orchestration & Resilience", font_size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4), "Four execution patterns, conversation memory, custom tools, and a 4-tier fallback that never fails", font_size=13, color=GRAY_500)

# ── Left: Orchestration Patterns ──
add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4), "Execution Patterns", font_size=16, color=GRAY_900, bold=True)

patterns = [
    ("Sequential Pipeline", "Task 1 → Task 2 → Task 3\nEach agent receives output from previous.\nUsed for: Blueprint generation, vendor selection",
     "Budget Agent → Vendor Agent → Timeline Agent", ROSE_500),
    ("Parallel + Sequential", "Phase 1: [Budget] ∥ [Style]  (parallel, independent)\nPhase 2: Vendor Agent  (uses Phase 1 outputs)\nPhase 3: Timeline Agent  (uses all outputs)",
     "ThreadPoolExecutor with 2 workers", BLUE_500),
    ("Self-Correction Loop", "Agent runs → Confidence check → Retry if < 60%\nPenalties: sparse output (-0.4), invalid JSON (-0.5),\nhallucination markers (-0.3), incomplete fields (-0.2)",
     "Max 1 retry with previous attempt injected", PURPLE_500),
    ("Two-Agent Collaboration", "Budget Agent (pricing analysis) →\nVendor Agent (competitive bid positioning)\nUsed for: bid_assistant, listing_insights",
     "Sequential crew, 2 tasks", EMERALD),
]
for i, (title, desc, detail, color) in enumerate(patterns):
    y = Inches(2.0 + i * 1.25)
    card = add_shape(slide, Inches(0.6), y, Inches(6.0), Inches(1.1), GRAY_50)
    add_rect(slide, Inches(0.6), y, Inches(0.06), Inches(1.1), color)
    # Pattern number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.15), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, Inches(0.88), y + Inches(0.17), Inches(0.3), Inches(0.3), str(i+1), font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.35), y + Inches(0.08), Inches(5.0), Inches(0.22), title, font_size=11, color=GRAY_900, bold=True)
    add_text(slide, Inches(1.35), y + Inches(0.32), Inches(5.0), Inches(0.7), desc, font_size=8, color=GRAY_600)

# ── Right: Fallback Architecture ──
add_text(slide, Inches(7.2), Inches(1.5), Inches(5), Inches(0.4), "4-Tier Resilience Architecture", font_size=16, color=GRAY_900, bold=True)

tiers = [
    ("Tier 1", "CrewAI Agents", "Multi-agent orchestration with\ncustom tools (query_vendors,\nget_market_rates, get_blueprint)\nOllama GLM4 local model", ROSE_500),
    ("Tier 2", "Gemini API", "Google Gemini 2.0 Flash\nFast cloud inference\nThinking tokens for complex tasks\nDirect API call via _ask_gemini()", BLUE_500),
    ("Tier 3", "Ollama Local", "Local LLM fallback\nWorks fully offline\nLower quality but always available\nDirect API call via _ask_ollama()", AMBER_500),
    ("Tier 4", "Deterministic", "Rule-based fallback\nNo AI needed — always works\nHardcoded market rates + templates\nVendor faq_data for pricing", EMERALD),
]
for i, (tier, name, desc, color) in enumerate(tiers):
    y = Inches(2.0 + i * 1.25)
    card = add_shape(slide, Inches(7.2), y, Inches(5.5), Inches(1.1), GRAY_50)
    add_rect(slide, Inches(7.2), y, Inches(5.5), Inches(0.06), color)
    # Tier badge
    badge = add_shape(slide, Inches(7.4), y + Inches(0.2), Inches(0.8), Inches(0.3), color)
    add_text(slide, Inches(7.45), y + Inches(0.22), Inches(0.7), Inches(0.25), tier, font_size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.35), y + Inches(0.18), Inches(3), Inches(0.25), name, font_size=12, color=GRAY_900, bold=True)
    add_text(slide, Inches(7.4), y + Inches(0.55), Inches(5.1), Inches(0.5), desc, font_size=8, color=GRAY_600)
    # Arrow to next tier
    if i < 3:
        arrow_y = y + Inches(1.1)
        add_text(slide, Inches(9.7), arrow_y - Inches(0.05), Inches(0.5), Inches(0.2), "↓", font_size=12, color=GRAY_400, align=PP_ALIGN.CENTER)

# ── Bottom: Memory + Tools ──
add_text(slide, Inches(0.8), Inches(6.8), Inches(4), Inches(0.25), "🧠 Conversation Memory", font_size=10, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(7.05), Inches(5), Inches(0.25), "Per-couple WeddingMemory (last 20 msgs) injected into every agent task description", font_size=8, color=GRAY_500)

add_text(slide, Inches(5.5), Inches(6.8), Inches(4), Inches(0.25), "🔧 Custom Agent Tools", font_size=10, color=GRAY_900, bold=True)
add_text(slide, Inches(5.5), Inches(7.05), Inches(5), Inches(0.25), "query_vendors(cat, city)  ·  get_market_rates(cat, city)  ·  get_blueprint_specs(id)  ·  SerperDevTool", font_size=8, color=GRAY_500)

add_text(slide, Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.25), "⚙️ Config", font_size=10, color=GRAY_900, bold=True)
add_text(slide, Inches(10.5), Inches(7.05), Inches(2.5), Inches(0.25), "allow_delegation=False · max_iter=1-3 · verbose=True", font_size=8, color=GRAY_500)

add_slide_number(slide, 21, TOTAL)

# ── Save ──
output = "/Users/snehaaiyer/Downloads/shehnai-local/Shehnai_Product_Walkthrough.pptx"
prs.save(output)
print(f"✅ Updated presentation saved: {output}")
print(f"   Now {TOTAL} slides (added 3 agentic AI slides)")
