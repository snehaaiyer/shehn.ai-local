#!/usr/bin/env python3
"""Generate Shehn.AI 6-slide crisp product deck — light theme."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Logo paths ──
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(BASE_DIR, "shehnai-logo.png")       # Main logo with text
ICON_PATH = os.path.join(BASE_DIR, "shehnai-avatar.png")      # App icon

# ── Brand Colors ──
ROSE_500   = RGBColor(0xE8, 0x6F, 0x5B)  # matches SVG #E86F5B
ROSE_600   = RGBColor(0xD4, 0x5B, 0x48)
ROSE_100   = RGBColor(0xFF, 0xE8, 0xE4)
ROSE_50    = RGBColor(0xFF, 0xF5, 0xF3)
SAGE_600   = RGBColor(0x2F, 0x4F, 0x46)  # matches SVG #2F4F46
SAGE_700   = RGBColor(0x25, 0x3E, 0x36)
SAGE_100   = RGBColor(0xE4, 0xEB, 0xE4)
SAGE_50    = RGBColor(0xF0, 0xF5, 0xF0)
GRAY_900   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_800   = RGBColor(0x2D, 0x2D, 0x3F)
GRAY_700   = RGBColor(0x37, 0x41, 0x51)
GRAY_600   = RGBColor(0x4B, 0x55, 0x63)
GRAY_500   = RGBColor(0x6B, 0x72, 0x80)
GRAY_400   = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_300   = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_200   = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100   = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_50    = RGBColor(0xF9, 0xFA, 0xFB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
AMBER_500  = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_600  = RGBColor(0xD9, 0x77, 0x06)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
EMERALD_700= RGBColor(0x04, 0x7D, 0x57)
BLUE_500   = RGBColor(0x3B, 0x82, 0xF6)
BLUE_600   = RGBColor(0x25, 0x63, 0xEB)
PURPLE_500 = RGBColor(0x8B, 0x5C, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 6

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ──
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.08
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=GRAY_900, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=GRAY_700, line_spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return txBox

def add_card(slide, left, top, width, height, fill_color=WHITE, accent_color=ROSE_500):
    shape = add_shape(slide, left, top, width, height, fill_color)
    # Top accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()
    return shape

def add_slide_number(slide, num):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3), f"{num}/{TOTAL_SLIDES}", font_size=9, color=GRAY_400, align=PP_ALIGN.RIGHT)
    add_logo(slide)

def add_logo(slide, left=Inches(11.3), top=Inches(0.2), height=Inches(0.45)):
    """Add Shehn.AI main logo (with text) to slide."""
    if os.path.exists(LOGO_PATH) and os.path.getsize(LOGO_PATH) > 100:
        try:
            slide.shapes.add_picture(LOGO_PATH, left, top, height=height)
        except Exception:
            pass

def add_icon(slide, left, top, height=Inches(0.5)):
    """Add the app icon (infinity symbol)."""
    if os.path.exists(ICON_PATH) and os.path.getsize(ICON_PATH) > 100:
        try:
            slide.shapes.add_picture(ICON_PATH, left, top, height=height)
        except Exception:
            pass

def add_stat_box(slide, left, top, number, label, accent=ROSE_500):
    box = add_card(slide, left, top, Inches(2.7), Inches(1.15), WHITE, accent)
    add_text(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.3), Inches(0.5), number, font_size=28, color=accent, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.3), Inches(0.3), label, font_size=10, color=GRAY_600)
    return box


# ══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Subtle decorative elements
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ROSE_50
circle.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4.5), Inches(4.5), Inches(4.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = SAGE_50
circle2.line.fill.background()

# Bottom sage accent bar
add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), SAGE_600)

# Logo (large, top-left)
add_logo(slide, left=Inches(0.8), top=Inches(1.2), height=Inches(0.9))

add_text(slide, Inches(0.8), Inches(2.5), Inches(9), Inches(1), "Shehn.AI", font_size=60, color=SAGE_600, bold=True)
add_text(slide, Inches(0.8), Inches(3.5), Inches(9), Inches(0.6), "AI-Powered Indian Wedding Planning Platform", font_size=24, color=GRAY_700, bold=True)
add_rect(slide, Inches(0.8), Inches(4.2), Inches(1.2), Inches(0.04), ROSE_500)
add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.5),
         "Reverse Marketplace  \u2022  AI Blueprints  \u2022  Multi-Agent Orchestration  \u2022  Sealed Bidding",
         font_size=14, color=GRAY_500)
add_text(slide, Inches(0.8), Inches(5.3), Inches(10), Inches(0.3), "Product Walkthrough  |  April 2026", font_size=11, color=GRAY_400)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — MARKET OPPORTUNITY (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRAY_50)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), "The Opportunity", font_size=32, color=SAGE_600, bold=True)
add_rect(slide, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.04), ROSE_500)

# Stats row
add_stat_box(slide, Inches(0.5), Inches(1.3), "\u20b96.5 Lakh Cr", "Indian Wedding Industry", ROSE_500)
add_stat_box(slide, Inches(3.5), Inches(1.3), "1 Crore+", "Weddings Per Year", AMBER_500)
add_stat_box(slide, Inches(6.5), Inches(1.3), "\u20b939.5L Avg", "Per-Wedding Spend", EMERALD)
add_stat_box(slide, Inches(9.5), Inches(1.3), "6 Categories", "Venue \u2022 Photo \u2022 Catering \u2022 Decor \u2022 Makeup \u2022 Music", BLUE_500)

# Pain points + solution (two columns on white cards)
# Left card: Pain
add_card(slide, Inches(0.5), Inches(2.8), Inches(5.8), Inches(3.8), WHITE, ROSE_500)
add_text(slide, Inches(0.8), Inches(2.95), Inches(5.2), Inches(0.4), "Pain Points", font_size=18, color=ROSE_500, bold=True)
add_multiline(slide, Inches(0.8), Inches(3.45), Inches(5.2), Inches(3.0), [
    "\u2022  Couples contact 30\u201350 vendors manually, no price transparency",
    "\u2022  Budget planning is guesswork \u2014 no per-category benchmarks",
    "\u2022  Vendors rely on referrals, no structured lead pipeline",
    "\u2022  No platform does AI planning + vendor matching + sealed bidding",
], font_size=12, color=GRAY_700)

# Right card: Solution
add_card(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.8), WHITE, EMERALD)
add_text(slide, Inches(7.1), Inches(2.95), Inches(5.2), Inches(0.4), "Shehn.AI Solution", font_size=18, color=EMERALD_700, bold=True)
add_multiline(slide, Inches(7.1), Inches(3.45), Inches(5.2), Inches(3.0), [
    "\u2022  AI generates complete wedding blueprint from preferences",
    "\u2022  Reverse marketplace: couple publishes, vendors bid",
    "\u2022  Sealed bids \u2014 couples see all quotes, vendors see bid count",
    "\u2022  Category-specific structured pricing for apples-to-apples comparison",
], font_size=12, color=GRAY_700)

add_text(slide, Inches(0.8), Inches(6.85), Inches(10), Inches(0.3),
         "Source: WedMeGood Annual Wedding Industry Report 2025-26", font_size=8, color=GRAY_400)
add_slide_number(slide, 2)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — HOW IT WORKS (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), "How It Works", font_size=32, color=SAGE_600, bold=True)
add_rect(slide, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.04), ROSE_500)

# Left: Couple Journey
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4), "Couple Journey", font_size=18, color=ROSE_500, bold=True)

couple_steps = [
    ("\u2776", "Set Preferences", "Date, city, budget, guest count, theme, events"),
    ("\u2777", "AI Blueprint", "CrewAI agents generate budget, vendor specs, timeline, tips"),
    ("\u2778", "Review & Publish", "Edit allocations, then publish to vendor marketplace"),
    ("\u2779", "Compare Quotes", "Side-by-side vendor quotes with budget comparison"),
    ("\u277A", "Book & Manage", "Message vendors, confirm bookings, manage RSVP + timeline"),
]
y = Inches(1.85)
for num, title, desc in couple_steps:
    add_card(slide, Inches(0.8), y, Inches(5.5), Inches(0.72), GRAY_50, ROSE_500)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(0.35), Inches(0.3), num, font_size=16, color=ROSE_500, bold=True)
    add_text(slide, Inches(1.45), y + Inches(0.1), Inches(4.5), Inches(0.28), title, font_size=13, color=GRAY_900, bold=True)
    add_text(slide, Inches(1.45), y + Inches(0.38), Inches(4.5), Inches(0.25), desc, font_size=10, color=GRAY_500)
    y += Inches(0.82)

# Right: Vendor Journey
add_text(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4), "Vendor Journey", font_size=18, color=SAGE_600, bold=True)

vendor_steps = [
    ("\u2776", "Register & Get Approved", "Business profile, portfolio, services \u2192 admin approval"),
    ("\u2777", "Browse Marketplace", "See published blueprints filtered to your category"),
    ("\u2778", "Submit Structured Quote", "Category-specific pricing \u2192 auto-calculated total"),
    ("\u2779", "Negotiate & Close", "In-app messaging, revise quotes, confirm booking"),
]
y = Inches(1.85)
for num, title, desc in vendor_steps:
    add_card(slide, Inches(7.0), y, Inches(5.5), Inches(0.72), GRAY_50, SAGE_600)
    add_text(slide, Inches(7.2), y + Inches(0.1), Inches(0.35), Inches(0.3), num, font_size=16, color=SAGE_600, bold=True)
    add_text(slide, Inches(7.65), y + Inches(0.1), Inches(4.5), Inches(0.28), title, font_size=13, color=GRAY_900, bold=True)
    add_text(slide, Inches(7.65), y + Inches(0.38), Inches(4.5), Inches(0.25), desc, font_size=10, color=GRAY_500)
    y += Inches(0.82)

# Bottom screens strip
add_rect(slide, Inches(0), Inches(6.3), SLIDE_W, Inches(1.2), SAGE_600)
add_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3), "Key Screens", font_size=11, color=WHITE, bold=True)
screens = ["Smart Planner", "Blueprint Review", "Vendor Marketplace", "Quote Dashboard", "In-App Messages", "RSVP Manager", "Budget Tracker"]
x = Inches(0.8)
for s in screens:
    add_shape(slide, x, Inches(6.75), Inches(1.55), Inches(0.45), SAGE_700)
    add_text(slide, x, Inches(6.82), Inches(1.55), Inches(0.3), s, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
    x += Inches(1.65)

add_slide_number(slide, 3)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — AGENTIC AI ARCHITECTURE (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GRAY_50)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), "Agentic AI: Multi-Agent Orchestration", font_size=32, color=SAGE_600, bold=True)
add_rect(slide, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.04), ROSE_500)

# 3 Crew boxes
crew_data = [
    ("Research Crew", ROSE_500, [
        "Venue Scout \u2014 location + capacity match",
        "Catering Analyst \u2014 cuisine + pricing",
        "Photo/Video Strategist \u2014 style + packages",
        "Decor Planner \u2014 theme + mandap design",
        "Entertainment Curator \u2014 DJ/band/sangeet",
    ]),
    ("Planning Crew", SAGE_600, [
        "Budget Optimizer \u2014 category allocation",
        "Timeline Architect \u2014 multi-day scheduling",
        "Vendor Matcher \u2014 preference scoring",
        "Guest Logistics \u2014 RSVP + seating",
        "Culture Specialist \u2014 ritual integration",
    ]),
    ("Content Crew", BLUE_600, [
        "Blueprint Writer \u2014 full plan document",
        "Cost Analyzer \u2014 savings identification",
        "Communication Drafter \u2014 vendor briefs",
        "Quality Reviewer \u2014 consistency check",
        "Summary Generator \u2014 couple output",
    ]),
]

x = Inches(0.5)
for crew_name, accent, agents in crew_data:
    add_card(slide, x, Inches(1.3), Inches(3.9), Inches(3.55), WHITE, accent)
    add_text(slide, x + Inches(0.25), Inches(1.48), Inches(3.4), Inches(0.35), crew_name, font_size=16, color=accent, bold=True)
    add_text(slide, x + Inches(0.25), Inches(1.82), Inches(3.4), Inches(0.22), "5 Specialized Agents", font_size=10, color=GRAY_400)
    y = Inches(2.15)
    for agent in agents:
        add_shape(slide, x + Inches(0.15), y, Inches(3.6), Inches(0.42), GRAY_50)
        add_text(slide, x + Inches(0.3), y + Inches(0.08), Inches(3.3), Inches(0.3), agent, font_size=9, color=GRAY_700)
        y += Inches(0.48)
    x += Inches(4.15)

# Bottom row: Orchestration + Fallback
add_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.03), GRAY_200)

# Orchestration patterns
add_text(slide, Inches(0.5), Inches(5.3), Inches(5.5), Inches(0.35), "Orchestration Patterns", font_size=14, color=SAGE_600, bold=True)
patterns = [
    "Sequential: Research \u2192 Planning \u2192 Content (blueprint gen)",
    "Parallel + Sequential: 5 research agents run concurrently",
    "Self-Correction: Quality reviewer triggers re-generation",
    "Two-Agent Collab: Budget Optimizer + Cost Analyzer iterate",
]
y = Inches(5.65)
for p in patterns:
    add_text(slide, Inches(0.5), y, Inches(6.0), Inches(0.22), f"\u25B8  {p}", font_size=9, color=GRAY_600)
    y += Inches(0.25)

# 4-tier fallback
add_text(slide, Inches(7.2), Inches(5.3), Inches(5.5), Inches(0.35), "4-Tier AI Fallback", font_size=14, color=SAGE_600, bold=True)

tiers = [
    ("CrewAI + Gemini", "Multi-agent orchestration (primary)", EMERALD),
    ("Gemini Direct", "Single-model generation (fallback 1)", BLUE_500),
    ("Ollama Local", "On-device LLM (fallback 2)", PURPLE_500),
    ("Deterministic", "Rule-based templates (guaranteed)", AMBER_500),
]
y = Inches(5.7)
for name, desc, color in tiers:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y + Inches(0.04), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, Inches(7.45), y, Inches(2.2), Inches(0.22), name, font_size=10, color=GRAY_900, bold=True)
    add_text(slide, Inches(9.7), y, Inches(3.0), Inches(0.22), desc, font_size=9, color=GRAY_500)
    y += Inches(0.3)

add_slide_number(slide, 4)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE & AUTOMATION (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), "Architecture & Automation", font_size=32, color=SAGE_600, bold=True)
add_rect(slide, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.04), ROSE_500)

# Left: Tech Stack
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4), "Tech Stack", font_size=18, color=ROSE_500, bold=True)

stack_items = [
    ("Frontend", "React 18 + TypeScript + Tailwind CSS + Framer Motion", ROSE_500),
    ("State", "Zustand with localStorage persistence", ROSE_500),
    ("Backend", "FastAPI (Python) + in-memory data stores", SAGE_600),
    ("AI Engine", "CrewAI + Google Gemini + Anthropic Claude Haiku", SAGE_600),
    ("Database", "NocoDB (Airtable-like, self-hosted)", BLUE_500),
    ("Vendor Data", "Serper API for real-time vendor search", BLUE_500),
]
y = Inches(1.8)
for label, desc, accent in stack_items:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.06), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    add_text(slide, Inches(1.15), y, Inches(1.5), Inches(0.25), label, font_size=11, color=GRAY_900, bold=True)
    add_text(slide, Inches(2.7), y, Inches(3.8), Inches(0.25), desc, font_size=10, color=GRAY_500)
    y += Inches(0.32)

# Data flow
add_rect(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(0.03), GRAY_200)
add_text(slide, Inches(0.8), Inches(4.05), Inches(5.5), Inches(0.35), "Data Flow", font_size=14, color=SAGE_600, bold=True)

flow_boxes = [
    ("React UI", ROSE_500, Inches(0.8)),
    ("FastAPI", SAGE_600, Inches(3.0)),
    ("CrewAI", BLUE_500, Inches(5.2)),
]
for label, color, left_pos in flow_boxes:
    add_card(slide, left_pos, Inches(4.5), Inches(1.6), Inches(0.55), WHITE, color)
    add_text(slide, left_pos, Inches(4.62), Inches(1.6), Inches(0.3), label, font_size=11, color=GRAY_900, bold=True, align=PP_ALIGN.CENTER)

# Arrows between flow boxes
add_text(slide, Inches(2.45), Inches(4.55), Inches(0.4), Inches(0.35), "\u2192", font_size=20, color=GRAY_400, bold=True)
add_text(slide, Inches(4.65), Inches(4.55), Inches(0.4), Inches(0.35), "\u2192", font_size=20, color=GRAY_400, bold=True)

add_text(slide, Inches(0.8), Inches(5.2), Inches(6), Inches(0.25),
         "React \u2192 Zustand Store \u2192 FastAPI REST \u2192 CrewAI Agents \u2192 Gemini/Claude \u2192 NocoDB",
         font_size=9, color=GRAY_400)

# Right: AI-Automated Features
add_text(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4), "AI-Automated Features", font_size=18, color=EMERALD, bold=True)

auto_features = [
    ("\u2713", "Wedding Blueprint Generation", "Complete plan from 5 inputs"),
    ("\u2713", "Budget Allocation", "Per-category splits \u2014 guest count + city"),
    ("\u2713", "Vendor Matching & Scoring", "Style / budget / location match"),
    ("\u2713", "Quote Total Estimation", "Unit pricing \u00d7 requirements auto-calc"),
    ("\u2713", "Timeline Generation", "Multi-day event scheduling"),
    ("\u2713", "Cost-Saving Recommendations", "AI-identified savings per category"),
    ("\u2713", "Smart Chat Planning", "Natural language wedding assistant"),
    ("\u2713", "RSVP Link Generation", "Unique invite codes with QR tracking"),
]
y = Inches(1.8)
for check, title, desc in auto_features:
    add_text(slide, Inches(7.0), y, Inches(0.3), Inches(0.25), check, font_size=12, color=EMERALD, bold=True)
    add_text(slide, Inches(7.35), y, Inches(2.5), Inches(0.25), title, font_size=11, color=GRAY_900, bold=True)
    add_text(slide, Inches(9.9), y, Inches(3.0), Inches(0.25), desc, font_size=9, color=GRAY_500)
    y += Inches(0.35)

# Role-based access
add_rect(slide, Inches(7.0), Inches(4.85), Inches(5.5), Inches(0.03), GRAY_200)
add_text(slide, Inches(7.0), Inches(5.05), Inches(5.5), Inches(0.3), "3 Role-Based Experiences", font_size=12, color=SAGE_600, bold=True)
roles = [
    ("Couple:", "13 screens \u2014 preferences, blueprint, quotes, messages, RSVP, budget, AI chat", ROSE_500),
    ("Vendor:", "5 screens \u2014 register, marketplace, quote form, inbox, profile", SAGE_600),
    ("Admin:", "2 screens \u2014 vendor approvals, user management", AMBER_500),
]
y = Inches(5.4)
for role, desc, color in roles:
    add_text(slide, Inches(7.0), y, Inches(0.8), Inches(0.22), role, font_size=10, color=color, bold=True)
    add_text(slide, Inches(7.7), y, Inches(4.5), Inches(0.22), desc, font_size=9, color=GRAY_500)
    y += Inches(0.25)

add_slide_number(slide, 5)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — CLOSING (Light)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Subtle decorative shapes
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4.5), Inches(4.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ROSE_50
circle.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(3.5), Inches(3.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = SAGE_50
circle2.line.fill.background()

# Logo centered
add_logo(slide, left=Inches(4.5), top=Inches(1.0), height=Inches(1.1))

add_text(slide, Inches(0), Inches(2.4), SLIDE_W, Inches(0.8), "Shehn.AI", font_size=48, color=SAGE_600, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.5), "AI-Powered Wedding Planning for Every Indian Celebration", font_size=18, color=GRAY_600, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.8), Inches(3.9), Inches(1.7), Inches(0.04), ROSE_500)

# Key metrics
metrics = [
    ("\u20b96.5L Cr", "Market Size"),
    ("15 AI Agents", "3 CrewAI Crews"),
    ("6 Categories", "Structured Pricing"),
    ("3 Roles", "Couple \u2022 Vendor \u2022 Admin"),
]
x = Inches(1.5)
for val, label in metrics:
    add_text(slide, x, Inches(4.3), Inches(2.5), Inches(0.4), val, font_size=22, color=SAGE_600, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.75), Inches(2.5), Inches(0.3), label, font_size=10, color=GRAY_500, align=PP_ALIGN.CENTER)
    x += Inches(2.7)

add_text(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.4),
         "Reverse Marketplace  \u2022  Agentic AI  \u2022  Sealed Bidding  \u2022  End-to-End Automation",
         font_size=13, color=GRAY_400, align=PP_ALIGN.CENTER)

# Bottom sage bar (matching slide 1)
add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), SAGE_600)
add_text(slide, Inches(0), Inches(7.08), SLIDE_W, Inches(0.3), "Built with CrewAI + Gemini + Claude + React + FastAPI",
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

add_slide_number(slide, 6)


# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
output_path = os.path.join(BASE_DIR, "Shehnai_Product_Walkthrough.pptx")
prs.save(output_path)
print(f"\n\u2705 Presentation saved: {output_path}")
print(f"   {TOTAL_SLIDES} slides, light theme, single logo")
